"""
Service Présentations — génère un diaporama (slides JSON) par IA + export PPTX
=============================================================================
À partir d'un groupe de documents (résumés / extraits), l'IA locale structure une
présentation en diapositives. La même structure « slides » alimente la visionneuse
(reveal.js) et l'export **PPTX** (python-pptx).
"""

import io
import json
import re

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from config import get_settings
from logger import get_logger
from models.document import Document
from models.metadata import MetadonneeIA
from services.ollama_service import OllamaService

log = get_logger(__name__)
settings = get_settings()

PROMPT_PRESENTATION = """Tu es un assistant qui crée des présentations professionnelles claires.
À partir des documents fournis, génère une présentation structurée en diapositives.
Réponds UNIQUEMENT par un objet JSON valide (sans markdown), de la forme :
{
  "titre": "Titre de la présentation",
  "theme": "un thème en 1-3 mots",
  "slides": [
    {"titre": "Titre de la diapo", "points": ["point clé 1", "point clé 2", "..."]}
  ]
}
Règles : 6 à 12 diapositives, 3 à 6 points par diapo, concis et factuel, en français.
Première diapo = introduction/contexte, dernière = synthèse/conclusion."""


def _extraire_json(texte: str) -> dict:
    """Extrait le premier objet JSON d'une réponse LLM."""
    try:
        return json.loads(texte)
    except json.JSONDecodeError:
        m = re.search(r"\{.*\}", texte, re.DOTALL)
        if m:
            return json.loads(m.group(0))
        raise


async def generer_slides(document_ids: list[str], db: AsyncSession, consigne: str | None = None,
                         model: str | None = None) -> dict:
    """Construit la structure de présentation via l'IA. Retourne {titre, theme, slides}."""
    import uuid as _uuid
    ids = [_uuid.UUID(i) for i in document_ids]
    rows = (await db.execute(
        select(Document, MetadonneeIA)
        .outerjoin(MetadonneeIA, MetadonneeIA.document_id == Document.id)
        .where(Document.id.in_(ids))
    )).all()

    if not rows:
        raise ValueError("Aucun document trouvé pour la présentation")

    # Contexte : titre + résumé (ou extrait tronqué) de chaque document
    blocs = []
    for doc, meta in rows:
        resume = (meta.resume if meta and meta.resume else (doc.texte_extrait or ""))[:1500]
        blocs.append(f"--- Document : {doc.nom} ---\n{resume}")
    contexte = "\n\n".join(blocs)[:14000]

    consigne_txt = f"\n\nConsigne de l'utilisateur : {consigne}" if consigne else ""
    prompt = f"{PROMPT_PRESENTATION}{consigne_txt}\n\nDocuments :\n{contexte}"

    ollama = OllamaService()
    from services import runtime_config
    used = model or runtime_config.model_for("rapport")
    reponse = await ollama.generate(prompt, model=used, format="json")
    data = _extraire_json(reponse)

    # Normalisation défensive
    slides = data.get("slides") or []
    slides = [
        {"titre": str(s.get("titre", "")).strip() or f"Diapositive {i + 1}",
         "points": [str(p).strip() for p in (s.get("points") or []) if str(p).strip()]}
        for i, s in enumerate(slides) if isinstance(s, dict)
    ]
    if not slides:
        raise ValueError("L'IA n'a pas produit de diapositives")

    return {
        "titre": str(data.get("titre", "Présentation")).strip() or "Présentation",
        "theme": (str(data.get("theme")).strip() if data.get("theme") else None),
        "slides": slides,
        "modele_utilise": used,
    }


def slides_to_pptx(titre: str, slides: list[dict]) -> bytes:
    """Construit un fichier PPTX à partir de la structure de slides."""
    from pptx import Presentation as Pptx
    from pptx.util import Pt

    prs = Pptx()

    # Diapo de titre
    titre_layout = prs.slide_layouts[0]
    s0 = prs.slides.add_slide(titre_layout)
    s0.shapes.title.text = titre
    if len(s0.placeholders) > 1:
        s0.placeholders[1].text = "Généré par Matothèque — IA locale"

    # Diapos contenu (titre + puces)
    contenu_layout = prs.slide_layouts[1]
    for sl in slides:
        s = prs.slides.add_slide(contenu_layout)
        s.shapes.title.text = sl.get("titre", "")
        body = s.placeholders[1].text_frame
        body.clear()
        points = sl.get("points") or []
        for i, p in enumerate(points):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = p
            para.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
